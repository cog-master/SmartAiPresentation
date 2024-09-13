from pptx import Presentation
from random import choice
import google.generativeai as genai
from dotenv import load_dotenv
from os import listdir, getenv
from json import loads

class generate_ppt:

    def __init__(self, topic: str) -> None:
        self.topic = topic
        path = f'templates/{choice(listdir("templates/"))}'
        self.presentation = Presentation(path)

    def get_data(self):

        load_dotenv()

        genai.configure(api_key=getenv('GEMINI_API_KEY'))
        model = genai.GenerativeModel('gemini-pro')

        prompt = prompt = f"""provide me 5 slide of the topic '{self.topic}' for my presentation in a json format like this: ```json {{"slide no": {{"slide title": "title", "slide content": "slide content single pararaph"}}}}``` consider the format of the json as shown in the example above."""

        while True:
            try: 
                response = model.generate_content(prompt)
                json_response = loads(response.text.replace('```','').replace('json',''))
                return json_response
            except: 
                print('Error Trying again...')

    def list_text_boxes(self, presentation, slide_num): # helper function to list all text boxes in a slide
        slide = presentation.slides[slide_num-1]
        text_boxes = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text:
                text_boxes.append(shape.text)
        return text_boxes

    def get_text_box_id(self, slide_no): # helper function to list all text boxes in a slide
        for idx, text in enumerate(self.list_text_boxes(Presentation('templates/professional_1.pptx'), slide_no), 1):
            print(f"slide no. {slide_no}: Text Box {idx}: {text} \n")

    def update_text_of_textbox(self, slide, text_box_id, new_text):
        slide = self.presentation.slides[(slide - 1)]
        count = 0
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text:
                count += 1
                if count == text_box_id:
                    text_frame = shape.text_frame
                    first_paragraph = text_frame.paragraphs[0]
                    first_run = first_paragraph.runs[0] if first_paragraph.runs else first_paragraph.add_run()
                    first_run.text = new_text
                    return
        raise ValueError(f"Slide {slide} or Text Box ID {text_box_id} not found")

    def update_ppxt(self):
        data = self.get_data()

        self.update_text_of_textbox(1, 1, data['slide 1']['slide title'] )   # slide 1
        self.update_text_of_textbox(1, 2, data['slide 1']['slide content'])  # slide 1

        self.update_text_of_textbox(2, 1, data['slide 2']['slide title'] )   # slide 2
        self.update_text_of_textbox(2, 2, data['slide 2']['slide content'])  # slide 2

        self.update_text_of_textbox(3, 1, data['slide 3']['slide title'] )   # slide 3
        self.update_text_of_textbox(3, 2, data['slide 3']['slide content'])  # slide 3

        self.update_text_of_textbox(4, 1, data['slide 4']['slide title'] )   # slide 4
        self.update_text_of_textbox(4, 2, data['slide 4']['slide content'])  # slide 4

        self.update_text_of_textbox(5, 1, data['slide 5']['slide title'] )   # slide 5
        self.update_text_of_textbox(5, 2, data['slide 5']['slide content'])  # slide 5
        
        file_name = f"{self.topic.replace(' ','_')}.pptx"
        self.presentation.save(f'{file_name}')
        return f'PPT saved: {file_name}'
